"""
RAG Server - FastAPI REST API for PDF Knowledge Base

Endpoints:
- GET  /health     - Health check
- POST /documents  - Add document (base64 content)
- GET  /documents  - List indexed files
- POST /query      - Search knowledge base
"""

from fastapi import FastAPI, HTTPException
from pydantic import BaseModel, Field
import os
import sys
import base64
import tempfile
import warnings

# ---------------------------------------------------------
# Set model cache directory to ./models/ (relative to script location)
# ---------------------------------------------------------

def get_base_path():
    """
    Determine execution environment:
    - If packaged exe, use executable folder (sys.executable)
    - If python script, use script folder (__file__)
    """
    if getattr(sys, 'frozen', False):
        return os.path.dirname(sys.executable)
    else:
        # Go up one level from server/ to project root
        return os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

BASE_DIR = get_base_path()
MODELS_DIR = os.path.join(BASE_DIR, "models")
DB_DIR = os.path.join(BASE_DIR, "faiss_index")

# Create directories if not exist
os.makedirs(MODELS_DIR, exist_ok=True)
os.makedirs(DB_DIR, exist_ok=True)

# Set environment variables for model cache BEFORE importing transformers
os.environ["HF_HOME"] = MODELS_DIR
os.environ["TRANSFORMERS_CACHE"] = MODELS_DIR
os.environ["SENTENCE_TRANSFORMERS_HOME"] = MODELS_DIR

# Suppress warnings
warnings.filterwarnings("ignore")

# ---------------------------------------------------------
# Constants
# ---------------------------------------------------------

MAX_FILE_SIZE_MB = 20
MAX_FILE_SIZE_BYTES = MAX_FILE_SIZE_MB * 1024 * 1024

SUPPORTED_EXTENSIONS = {
    '.pdf', '.docx', '.pptx', '.xlsx', '.xls',
    '.md', '.markdown',
    '.txt', '.log',
    '.bat', '.sh', '.ps1',
    '.json', '.yaml', '.yml', '.ini', '.cfg', '.conf',
    '.csv',
    '.py', '.js', '.ts', '.html', '.css', '.xml'
}

TEXT_EXTENSIONS = {
    '.txt', '.log',
    '.bat', '.sh', '.ps1',
    '.json', '.yaml', '.yml', '.ini', '.cfg', '.conf',
    '.csv',
    '.py', '.js', '.ts', '.html', '.css', '.xml'
}

# ---------------------------------------------------------
# Lazy loading with Singleton pattern
# ---------------------------------------------------------

_embedding_function = None
_db_cache = None

def get_embedding_function():
    """Lazy load embedding model."""
    global _embedding_function
    if _embedding_function is None:
        from langchain_huggingface import HuggingFaceEmbeddings
        
        print("[INFO] Loading AI model (first run may be slow)...", file=sys.stderr)
        _embedding_function = HuggingFaceEmbeddings(
            model_name="jinaai/jina-embeddings-v3",
            model_kwargs={"trust_remote_code": True}
        )
        print("[OK] Model loaded successfully", file=sys.stderr)
    return _embedding_function

def get_db():
    """Get DB instance with Singleton pattern."""
    global _db_cache
    
    if _db_cache is not None:
        return _db_cache

    from langchain_community.vectorstores import FAISS
    
    embedding_func = get_embedding_function()
    
    if os.path.exists(DB_DIR) and os.path.exists(os.path.join(DB_DIR, "index.faiss")):
        print("[INFO] Loading FAISS index from disk...", file=sys.stderr)
        _db_cache = FAISS.load_local(DB_DIR, embedding_func, allow_dangerous_deserialization=True)
        return _db_cache
        
    return None

def save_db(db):
    """Save index to disk and update global cache."""
    global _db_cache
    db.save_local(DB_DIR)
    _db_cache = db

# ---------------------------------------------------------
# Document Loading (reused from original server.py)
# ---------------------------------------------------------

def get_file_extension(file_path: str) -> str:
    """Get file extension (lowercase)."""
    return os.path.splitext(file_path)[1].lower()

def load_document(file_path: str):
    """
    Universal document loader: auto-selects appropriate loader based on file type.
    
    Returns:
        list: List of Document objects
    """
    from langchain_core.documents import Document
    
    ext = get_file_extension(file_path)
    file_name = os.path.basename(file_path)
    
    try:
        if ext == '.pdf':
            try:
                import pypdf
            except ImportError:
                print("[ERROR] Missing pypdf module", file=sys.stderr)
                return []
            from langchain_community.document_loaders import PyPDFLoader
            loader = PyPDFLoader(file_path)
            return loader.load()
        
        elif ext == '.docx':
            import docx2txt
            text = docx2txt.process(file_path)
            if text.strip():
                return [Document(
                    page_content=text,
                    metadata={"source": file_name, "file_type": "docx"}
                )]
            return []
        
        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(file_path)
            documents = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if slide_text:
                    documents.append(Document(
                        page_content="\n".join(slide_text),
                        metadata={"source": file_name, "page": slide_num, "file_type": "pptx"}
                    ))
            return documents
        
        elif ext in ['.xlsx', '.xls']:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            documents = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows_text = []
                for row in sheet.iter_rows(max_row=1000):
                    row_values = [str(cell.value) if cell.value is not None else "" for cell in row]
                    if any(v.strip() for v in row_values):
                        rows_text.append(" | ".join(row_values))
                if rows_text:
                    documents.append(Document(
                        page_content="\n".join(rows_text),
                        metadata={"source": file_name, "sheet": sheet_name, "file_type": "excel"}
                    ))
            return documents
        
        elif ext in ['.md', '.markdown']:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            if text.strip():
                return [Document(
                    page_content=text,
                    metadata={"source": file_name, "file_type": "markdown"}
                )]
            return []
        
        elif ext in TEXT_EXTENSIONS:
            import chardet
            
            try:
                with open(file_path, 'rb') as f:
                    raw_data = f.read()
                
                detected = chardet.detect(raw_data)
                encoding = detected.get('encoding', 'utf-8') or 'utf-8'
                confidence = detected.get('confidence', 0)
                
                text = raw_data.decode(encoding, errors='ignore')
                
                if text.strip():
                    file_type_map = {
                        '.txt': 'text', '.log': 'log',
                        '.bat': 'script', '.sh': 'script', '.ps1': 'script',
                        '.json': 'config', '.yaml': 'config', '.yml': 'config',
                        '.ini': 'config', '.cfg': 'config', '.conf': 'config',
                        '.csv': 'data',
                        '.py': 'code', '.js': 'code', '.ts': 'code',
                        '.html': 'code', '.css': 'code', '.xml': 'code'
                    }
                    file_type = file_type_map.get(ext, 'text')
                    
                    return [Document(
                        page_content=text,
                        metadata={
                            "source": file_name,
                            "file_type": file_type,
                            "encoding": encoding,
                            "encoding_confidence": round(confidence, 2)
                        }
                    )]
            except Exception as e:
                print(f"[ERROR] Reading text file {file_path}: {e}", file=sys.stderr)
            return []
        
        else:
            return []
    
    except Exception as e:
        print(f"[ERROR] Loading document {file_path}: {e}", file=sys.stderr)
        return []

# ---------------------------------------------------------
# FastAPI App
# ---------------------------------------------------------

app = FastAPI(
    title="RAG Server",
    description="REST API for PDF Knowledge Base",
    version="1.0.0"
)

# ---------------------------------------------------------
# Request/Response Models
# ---------------------------------------------------------

class AddDocumentRequest(BaseModel):
    filename: str = Field(..., description="Original filename with extension")
    content_base64: str = Field(..., description="Base64 encoded file content")

class AddDocumentResponse(BaseModel):
    status: str
    message: str
    chunks_added: int = 0

class QueryRequest(BaseModel):
    query: str = Field(..., description="Search query")

class QueryResult(BaseModel):
    source: str
    page: str
    content: str

class QueryResponse(BaseModel):
    query: str
    results: list[QueryResult]

class FileInfo(BaseModel):
    filename: str
    chunks: int
    pages: int | None = None

class ListDocumentsResponse(BaseModel):
    total_files: int
    total_chunks: int
    files: list[FileInfo]

class HealthResponse(BaseModel):
    status: str
    model_loaded: bool
    index_loaded: bool

# ---------------------------------------------------------
# Endpoints
# ---------------------------------------------------------

@app.get("/health", response_model=HealthResponse)
def health_check():
    """Health check endpoint."""
    return HealthResponse(
        status="ok",
        model_loaded=_embedding_function is not None,
        index_loaded=_db_cache is not None
    )

@app.post("/documents", response_model=AddDocumentResponse)
def add_document(request: AddDocumentRequest):
    """Add a document to the knowledge base."""
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    from langchain_community.vectorstores import FAISS
    
    # Validate filename extension
    ext = get_file_extension(request.filename)
    if ext not in SUPPORTED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported format '{ext}'. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )
    
    # Decode base64 content
    try:
        file_content = base64.b64decode(request.content_base64)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Invalid base64 content: {str(e)}")
    
    # Check file size
    if len(file_content) > MAX_FILE_SIZE_BYTES:
        raise HTTPException(
            status_code=400,
            detail=f"File exceeds {MAX_FILE_SIZE_MB}MB limit (got {len(file_content) / 1024 / 1024:.1f}MB)"
        )
    
    # Write to temp file for processing
    try:
        with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp_file:
            tmp_file.write(file_content)
            tmp_path = tmp_file.name
        
        # Load and process document
        docs = load_document(tmp_path)
        
        if not docs:
            raise HTTPException(status_code=400, detail="Empty or unreadable document")
        
        # Split into chunks
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=600, chunk_overlap=100)
        splits = text_splitter.split_documents(docs)
        
        if not splits:
            raise HTTPException(status_code=400, detail="Document content is empty or unreadable")
        
        # Update metadata with original filename
        for split in splits:
            split.metadata["source"] = request.filename
        
        # Add to FAISS
        current_db = get_db()
        embedding_func = get_embedding_function()
        
        if current_db:
            current_db.add_documents(splits)
            save_db(current_db)
        else:
            new_db = FAISS.from_documents(splits, embedding_func)
            save_db(new_db)
        
        return AddDocumentResponse(
            status="ok",
            message=f"Added '{request.filename}' with {len(splits)} chunks to knowledge base.",
            chunks_added=len(splits)
        )
    
    finally:
        # Cleanup temp file
        if 'tmp_path' in locals() and os.path.exists(tmp_path):
            os.unlink(tmp_path)

@app.get("/documents", response_model=ListDocumentsResponse)
def list_documents():
    """List all indexed documents."""
    db = get_db()
    
    if not db:
        return ListDocumentsResponse(total_files=0, total_chunks=0, files=[])
    
    sources_info = {}
    total_chunks = 0
    
    for doc_id in db.docstore._dict:
        doc = db.docstore._dict[doc_id]
        source = doc.metadata.get('source', 'Unknown')
        page = doc.metadata.get('page', None)
        
        if source not in sources_info:
            sources_info[source] = {'pages': set(), 'chunks': 0}
        
        sources_info[source]['chunks'] += 1
        if page is not None:
            sources_info[source]['pages'].add(page)
        total_chunks += 1
    
    files = [
        FileInfo(
            filename=source,
            chunks=info['chunks'],
            pages=len(info['pages']) if info['pages'] else None
        )
        for source, info in sorted(sources_info.items())
    ]
    
    return ListDocumentsResponse(
        total_files=len(files),
        total_chunks=total_chunks,
        files=files
    )

@app.post("/query", response_model=QueryResponse)
def query_knowledge_base(request: QueryRequest):
    """Search the knowledge base."""
    db = get_db()
    
    if not db:
        raise HTTPException(status_code=404, detail="Knowledge base is empty. Please add documents first.")
    
    results = db.similarity_search(request.query, k=4)
    
    return QueryResponse(
        query=request.query,
        results=[
            QueryResult(
                source=doc.metadata.get('source', 'Unknown'),
                page=str(doc.metadata.get('page', 'N/A')),
                content=doc.page_content
            )
            for doc in results
        ]
    )

# ---------------------------------------------------------
# Startup Event
# ---------------------------------------------------------

@app.on_event("startup")
async def startup_event():
    """Pre-load models and index on startup."""
    print("[System] Initializing RAG Server...", file=sys.stderr)
    
    print("[System] 1/2 Loading Embedding Model...", file=sys.stderr)
    get_embedding_function()
    
    print("[System] 2/2 Pre-loading Database...", file=sys.stderr)
    get_db()
    
    print("[System] All systems ready.", file=sys.stderr)

# ---------------------------------------------------------
# Run with: uvicorn rag_server:app --host 0.0.0.0 --port 8000
# ---------------------------------------------------------

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
