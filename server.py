from fastmcp import FastMCP
import os
import sys
import warnings
import traceback

# ---------------------------------------------------------
# 修改點 1: 移除頂層的 Heavy Imports
# 只保留 os, sys, fastmcp 等輕量級套件
# 這樣 Server 啟動時就不會卡在載入 PyTorch/LangChain
# ---------------------------------------------------------

# 忽略警告
warnings.filterwarnings("ignore")

mcp = FastMCP("My Local Library")

def get_base_path():
    """
    判斷執行環境：
    如果是打包後的 exe，使用執行檔所在的資料夾 (sys.executable)
    如果是原始 python script，使用檔案所在的資料夾 (__file__)
    """
    if getattr(sys, 'frozen', False):
        # 如果是打包後的 exe
        return os.path.dirname(sys.executable)
    else:
        # 如果是開發中的 .py
        return os.path.dirname(os.path.abspath(__file__))

BASE_DIR = get_base_path()
DB_DIR = os.path.join(BASE_DIR, "faiss_index")

# ---------------------------------------------------------
# 修改點 2: 使用全域變數配合 Singleton 模式延遲載入
# ---------------------------------------------------------
_embedding_function = None

def get_embedding_function():
    """
    延遲載入 Embedding 模型。
    只有在第一次呼叫工具時才會執行 import 和模型下載。
    """
    global _embedding_function
    if _embedding_function is None:
        # 將 import 移到這裡，避免啟動時卡住
        from langchain_huggingface import HuggingFaceEmbeddings
        
        print("正在載入 AI 模型 (第一次執行會較慢)...", file=sys.stderr)
        _embedding_function = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
    return _embedding_function

def get_db():
    """統一取得 DB 實例 (包含延遲載入 FAISS)"""
    # 將 import 移到這裡
    from langchain_community.vectorstores import FAISS
    
    # 確保模型已載入
    embedding_func = get_embedding_function()
    
    if os.path.exists(DB_DIR) and os.path.exists(os.path.join(DB_DIR, "index.faiss")):
        return FAISS.load_local(DB_DIR, embedding_func, allow_dangerous_deserialization=True)
    return None

def save_db(db):
    """將索引存回硬碟"""
    db.save_local(DB_DIR)

# 支援的檔案格式
SUPPORTED_EXTENSIONS = {
    # 文件類
    '.pdf', '.docx', '.pptx', '.xlsx', '.xls',
    # Markdown
    '.md', '.markdown',
    # 純文字
    '.txt', '.log',
    # 腳本
    '.bat', '.sh', '.ps1',
    # 設定檔
    '.json', '.yaml', '.yml', '.ini', '.cfg', '.conf',
    # 資料檔
    '.csv',
    # 程式碼
    '.py', '.js', '.ts', '.html', '.css', '.xml'
}

# 純文字類型副檔名（用於統一處理）
TEXT_EXTENSIONS = {
    '.txt', '.log',
    '.bat', '.sh', '.ps1',
    '.json', '.yaml', '.yml', '.ini', '.cfg', '.conf',
    '.csv',
    '.py', '.js', '.ts', '.html', '.css', '.xml'
}

# 批次處理設定 - 每處理 N 個檔案就寫入一次 FAISS，降低記憶體使用並增加可靠性
BATCH_SIZE = 10

def get_indexed_sources(db) -> set:
    """
    取得已索引的檔案名稱集合，用於重複檢測。
    
    Returns:
        set: 已索引的檔案名稱（source）集合
    """
    if not db or not hasattr(db, 'docstore') or not hasattr(db.docstore, '_dict'):
        return set()
    return {doc.metadata.get('source') for doc in db.docstore._dict.values() if doc.metadata.get('source')}

def get_file_extension(file_path: str) -> str:
    """取得檔案副檔名（小寫）"""
    return os.path.splitext(file_path)[1].lower()

def load_document(file_path: str):
    """
    通用文件載入器：根據檔案類型自動選擇適合的 loader
    
    支援格式：
    - PDF (.pdf)
    - Word (.docx)
    - PowerPoint (.pptx)
    - Excel (.xlsx, .xls)
    - Markdown (.md, .markdown)
    - 純文字 (.txt, .log)
    - 腳本 (.bat, .sh, .ps1)
    - 設定檔 (.json, .yaml, .yml, .ini, .cfg, .conf)
    - 資料檔 (.csv)
    - 程式碼 (.py, .js, .ts, .html, .css, .xml)
    
    Returns:
        list: Document 物件列表
    """
    from langchain_core.documents import Document
    
    ext = get_file_extension(file_path)
    file_name = os.path.basename(file_path)
    
    try:
        if ext == '.pdf':
            try:
                import pypdf
            except ImportError:
                print("錯誤：缺少 pypdf 模組，請執行 pip install pypdf", file=sys.stderr)
                return []
            from langchain_community.document_loaders import PyPDFLoader
            loader = PyPDFLoader(file_path)
            return loader.load()
        
        elif ext == '.docx':
            import docx2txt
            text = docx2txt.process(file_path)
            if text.strip():
                return [Document(
                    page_content=text,
                    metadata={"source": file_name, "file_type": "docx"}
                )]
            return []
        
        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(file_path)
            documents = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if slide_text:
                    documents.append(Document(
                        page_content="\n".join(slide_text),
                        metadata={"source": file_name, "page": slide_num, "file_type": "pptx"}
                    ))
            return documents
        
        elif ext in ['.xlsx', '.xls']:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            documents = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows_text = []
                for row in sheet.iter_rows(max_row=1000):  # 限制最大行數避免記憶體問題
                    row_values = [str(cell.value) if cell.value is not None else "" for cell in row]
                    if any(v.strip() for v in row_values):
                        rows_text.append(" | ".join(row_values))
                if rows_text:
                    documents.append(Document(
                        page_content="\n".join(rows_text),
                        metadata={"source": file_name, "sheet": sheet_name, "file_type": "excel"}
                    ))
            return documents
        
        elif ext in ['.md', '.markdown']:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            if text.strip():
                return [Document(
                    page_content=text,
                    metadata={"source": file_name, "file_type": "markdown"}
                )]
            return []
        
        elif ext in TEXT_EXTENSIONS:
            # 使用 chardet 自動偵測編碼
            import chardet
            
            try:
                with open(file_path, 'rb') as f:
                    raw_data = f.read()
                
                # 偵測編碼
                detected = chardet.detect(raw_data)
                encoding = detected.get('encoding', 'utf-8') or 'utf-8'
                confidence = detected.get('confidence', 0)
                
                # 解碼文字
                text = raw_data.decode(encoding, errors='ignore')
                
                if text.strip():
                    # 根據副檔名決定 file_type
                    file_type_map = {
                        '.txt': 'text', '.log': 'log',
                        '.bat': 'script', '.sh': 'script', '.ps1': 'script',
                        '.json': 'config', '.yaml': 'config', '.yml': 'config',
                        '.ini': 'config', '.cfg': 'config', '.conf': 'config',
                        '.csv': 'data',
                        '.py': 'code', '.js': 'code', '.ts': 'code',
                        '.html': 'code', '.css': 'code', '.xml': 'code'
                    }
                    file_type = file_type_map.get(ext, 'text')
                    
                    return [Document(
                        page_content=text,
                        metadata={
                            "source": file_name,
                            "file_type": file_type,
                            "encoding": encoding,
                            "encoding_confidence": round(confidence, 2)
                        }
                    )]
            except Exception as e:
                print(f"讀取純文字檔案 {file_path} 時發生錯誤: {e}", file=sys.stderr)
            return []
        
        else:
            return []
    
    except Exception as e:
        print(f"載入文件 {file_path} 時發生錯誤: {e}", file=sys.stderr)
        return []

# ---------------------------------------------------------
# 修改點 3: 將工具所需的 Import 移至函式內部
# ---------------------------------------------------------

@mcp.tool()
def add_folder_to_library(folder_path: str):
    """[批次處理] 讀取資料夾內所有支援的文件並加入知識庫
    
    支援格式：PDF, DOCX, PPTX, XLSX, XLS, MD, TXT, LOG, BAT, SH, PS1, JSON, YAML, YML, INI, CFG, CONF, CSV, PY, JS, TS, HTML, CSS, XML
    
    優化功能：
    - 分批寫入 FAISS（每 N 個檔案寫入一次，降低記憶體使用）
    - 自動跳過已索引的檔案（重複檢測）
    - 進度回報（在 stderr 輸出處理進度）
    - 斷點續傳友善（分批寫入，中途失敗也保留部分成果）
    """
    # Import moved locally
    import glob
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    from langchain_community.vectorstores import FAISS

    folder_path = folder_path.strip('"').strip("'")
    
    if not os.path.exists(folder_path):
        return f"錯誤：找不到資料夾 -> {folder_path}"

    # 搜尋所有支援的檔案格式
    all_files = []
    for ext in SUPPORTED_EXTENSIONS:
        all_files.extend(glob.glob(os.path.join(folder_path, f"*{ext}")))
        # 也搜尋大寫副檔名
        all_files.extend(glob.glob(os.path.join(folder_path, f"*{ext.upper()}")))
    
    # 去重複並排序（確保處理順序一致）
    all_files = sorted(set(all_files))
    
    if not all_files:
        supported_list = ", ".join(SUPPORTED_EXTENSIONS)
        return f"在 '{folder_path}' 中找不到任何支援的檔案。\n支援格式: {supported_list}"

    # 取得已索引的檔案，用於重複檢測
    current_db = get_db()
    indexed_sources = get_indexed_sources(current_db)
    
    # 過濾掉已索引的檔案
    files_to_process = []
    skipped_files = []
    for file_path in all_files:
        file_name = os.path.basename(file_path)
        if file_name in indexed_sources:
            skipped_files.append(file_name)
        else:
            files_to_process.append(file_path)
    
    if not files_to_process:
        if skipped_files:
            return f"📋 資料夾中的 {len(skipped_files)} 個檔案都已在知識庫中，無需重複索引。"
        return "沒有找到需要處理的檔案。"

    # 進度回報
    total_files = len(files_to_process)
    print(f"[開始處理] 共 {total_files} 個新檔案待處理，已跳過 {len(skipped_files)} 個重複檔案", file=sys.stderr)
    
    processed_files = []
    failed_files = []
    total_splits_count = 0
    batch_count = 0
    
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=600, chunk_overlap=100)
    embedding_func = get_embedding_function()
    
    batch_splits = []  # 當前批次的片段

    for idx, file_path in enumerate(files_to_process, 1):
        file_name = os.path.basename(file_path)
        
        # 進度回報
        print(f"[{idx}/{total_files}] 處理中: {file_name}", file=sys.stderr)
        
        try:
            docs = load_document(file_path)
            if docs:
                splits = text_splitter.split_documents(docs)
                if splits:
                    for split in splits:
                        split.metadata["source"] = file_name
                    batch_splits.extend(splits)
                    processed_files.append(file_name)
                else:
                    failed_files.append((file_name, "文件內容為空"))
            else:
                failed_files.append((file_name, "無法讀取內容"))
        except Exception as e:
            failed_files.append((file_name, str(e)))
            print(f"[錯誤] {file_name}: {e}", file=sys.stderr)
        
        # 分批寫入 FAISS - 每 BATCH_SIZE 個檔案或最後一個檔案時寫入
        if len(processed_files) > 0 and (len(processed_files) % BATCH_SIZE == 0 or idx == total_files):
            if batch_splits:
                try:
                    batch_count += 1
                    print(f"[寫入批次 {batch_count}] 正在寫入 {len(batch_splits)} 個片段到 FAISS...", file=sys.stderr)
                    
                    # 重新取得最新的 DB（可能在上一批次已更新）
                    current_db = get_db()
                    
                    if current_db:
                        current_db.add_documents(batch_splits)
                        save_db(current_db)
                    else:
                        new_db = FAISS.from_documents(batch_splits, embedding_func)
                        save_db(new_db)
                    
                    total_splits_count += len(batch_splits)
                    batch_splits = []  # 清空批次，準備下一批
                    print(f"[寫入批次 {batch_count}] 完成！累計已寫入 {total_splits_count} 個片段", file=sys.stderr)
                    
                except Exception as e:
                    error_msg = f"寫入批次 {batch_count} 時發生錯誤: {str(e)}"
                    print(f"[嚴重錯誤] {error_msg}", file=sys.stderr)
                    # 記錄這批檔案為失敗（但保留之前批次的成果）
                    return f"⚠️ 部分處理完成，但在批次 {batch_count} 時發生錯誤。\n" \
                           f"📁 已成功處理: {len(processed_files) - len(batch_splits)} 個檔案\n" \
                           f"📄 已寫入: {total_splits_count} 個片段\n" \
                           f"❌ 錯誤: {error_msg}"

    # 組合最終結果
    result = f"✅ 批次處理完成！\n"
    result += f"{'='*40}\n"
    result += f"📁 共處理 {len(processed_files)} 個檔案\n"
    result += f"📄 新增 {total_splits_count} 個片段\n"
    result += f"📦 分 {batch_count} 個批次寫入\n"
    
    if skipped_files:
        result += f"\n⏭️ 已跳過 {len(skipped_files)} 個重複檔案（已存在於知識庫）\n"
    
    if failed_files:
        result += f"\n⚠️ 以下 {len(failed_files)} 個檔案處理失敗:\n"
        for file_name, reason in failed_files:
            result += f"  - {file_name}: {reason}\n"
    
    return result

@mcp.tool()
def add_pdf_to_library(pdf_path: str):
    """[單檔處理] 將 PDF 加入知識庫（向下相容，建議使用 add_document_to_library）"""
    return add_document_to_library(pdf_path)

@mcp.tool()
def add_document_to_library(file_path: str):
    """[單檔處理] 將文件加入知識庫
    
    支援格式：PDF, DOCX, PPTX, XLSX, XLS, MD, TXT, LOG, BAT, SH, PS1, JSON, YAML, YML, INI, CFG, CONF, CSV, PY, JS, TS, HTML, CSS, XML
    """
    # Import moved locally
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    from langchain_community.vectorstores import FAISS

    file_path = file_path.strip('"').strip("'")
    
    if not os.path.exists(file_path):
        return f"錯誤：找不到檔案 -> {file_path}"
    
    # 檢查檔案格式
    ext = get_file_extension(file_path)
    if ext not in SUPPORTED_EXTENSIONS:
        supported_list = ", ".join(SUPPORTED_EXTENSIONS)
        return f"錯誤：不支援的檔案格式 '{ext}'\n支援格式: {supported_list}"

    try:
        docs = load_document(file_path)
        
        if not docs:
            return f"文件內容為空或無法讀取: {os.path.basename(file_path)}"
        
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=600, chunk_overlap=100)
        splits = text_splitter.split_documents(docs)
        
        if not splits:
            return "文件內容為空或無法讀取。"

        # 確保 source metadata 正確
        for split in splits:
            split.metadata["source"] = os.path.basename(file_path)

        # 寫入 FAISS
        current_db = get_db()
        embedding_func = get_embedding_function()

        if current_db:
            current_db.add_documents(splits)
            save_db(current_db)
        else:
            new_db = FAISS.from_documents(splits, embedding_func)
            save_db(new_db)
        
        file_type_emoji = {
            '.pdf': '📕',
            '.docx': '📘',
            '.pptx': '📙',
            '.xlsx': '📗',
            '.xls': '📗',
            '.md': '📝',
            '.markdown': '📝',
            # 純文字
            '.txt': '📄',
            '.log': '📋',
            # 腳本
            '.bat': '⚙️',
            '.sh': '⚙️',
            '.ps1': '⚙️',
            # 設定檔
            '.json': '🔧',
            '.yaml': '🔧',
            '.yml': '🔧',
            '.ini': '🔧',
            '.cfg': '🔧',
            '.conf': '🔧',
            # 資料檔
            '.csv': '📊',
            # 程式碼
            '.py': '🐍',
            '.js': '💻',
            '.ts': '💻',
            '.html': '🌐',
            '.css': '🎨',
            '.xml': '📰'
        }
        emoji = file_type_emoji.get(ext, '📄')
        
        return f"✅ 成功！{emoji} 已將 '{os.path.basename(file_path)}' 的 {len(splits)} 個片段加入知識庫。"
    
    except Exception as e:
        return f"處理文件時發生錯誤: {str(e)}"

@mcp.tool()
def list_indexed_files():
    """[查詢] 列出知識庫中已索引的所有文件及統計資訊"""
    try:
        db = get_db()
        if not db:
            return "📭 知識庫目前是空的，請先使用 add_pdf_to_library 或 add_folder_to_library 加入文件。"

        # 從 docstore 中提取所有文件的 metadata
        sources_info = {}  # source -> {'pages': set(), 'chunks': count}
        total_chunks = 0
        
        for doc_id in db.docstore._dict:
            doc = db.docstore._dict[doc_id]
            source = doc.metadata.get('source', '未知來源')
            page = doc.metadata.get('page', None)
            
            if source not in sources_info:
                sources_info[source] = {'pages': set(), 'chunks': 0}
            
            sources_info[source]['chunks'] += 1
            if page is not None:
                sources_info[source]['pages'].add(page)
            total_chunks += 1

        # 組合輸出結果
        file_count = len(sources_info)
        result = f"📚 知識庫統計資訊\n"
        result += f"{'='*40}\n"
        result += f"📁 已索引文件數量: {file_count}\n"
        result += f"📄 總片段數量: {total_chunks}\n"
        result += f"{'='*40}\n\n"
        result += f"📋 已索引文件清單:\n"
        result += f"{'-'*40}\n"
        
        for idx, (source, info) in enumerate(sorted(sources_info.items()), 1):
            page_info = f", 共 {len(info['pages'])} 頁" if info['pages'] else ""
            result += f"{idx}. {source}\n"
            result += f"   └─ {info['chunks']} 個片段{page_info}\n"
        
        return result

    except Exception as e:
        return f"查詢索引時發生錯誤: {str(e)}"

@mcp.tool()
def query_library(query: str):
    """[搜尋] 從知識庫搜尋相關資訊"""
    try:
        # get_db 內部會處理 FAISS 的 import
        db = get_db()
        if not db:
            return "知識庫目前是空的，請先加入 PDF 檔案。"

        results = db.similarity_search(query, k=4)
        
        if not results:
            return "在知識庫中找不到相關資訊。"

        response_text = f"針對 '{query}' 的搜尋結果:\n\n"
        for doc in results:
            source = doc.metadata.get('source', '未知來源')
            page = doc.metadata.get('page', 'N/A')
            response_text += f"--- {source} (P.{page}) ---\n{doc.page_content}\n\n"
            
        return response_text
        
    except Exception as e:
        return f"搜尋時發生錯誤: {str(e)}"

if __name__ == "__main__":
    log_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), "debug_error.log")
    try:
        mcp.run()
    except Exception:
        with open(log_file, "w", encoding="utf-8") as f:
            f.write(traceback.format_exc())
        sys.exit(1)